from pptx import Presentation

from modules.core import SlidesgenPresentation, TitleSlide, AgendaSlide, SimpleContentSlide, TwoColumnsSlide, QuoteSlide, ImpressionSlide, SectionTransitionSlide, ThankYouSlide

# ==
def update_text_content_of_shape(shape, text: str):
    if shape.has_text_frame:
        text_frame = shape.text_frame
        
        if len(text_frame.paragraphs) == 0:
            p = text_frame.add_paragraph()
            
        else:
            p = text_frame.paragraphs[0]

        while p.runs:
            p._element.remove(p.runs[0]._element)

        run = p.add_run()
        run.text = text

def set_formatted_content(shape, content):
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame

    # Clear all paragraphs
    for _ in range(len(text_frame.paragraphs)):
        p = text_frame.paragraphs[0]
        text_frame._element.remove(p._element)

    lines = content.strip().split("\n")
    for line in lines:
        if line.startswith("-"):
            p = text_frame.add_paragraph()
            p.text = line[1:].strip()
            p.level = 0
            p.font.bold = None  # use template format
            p.font.size = None
            p.space_after = 0
            p.space_before = 0
            p.alignment = None
            p._element.set("marL", "0")  # Optional: remove indent
            
        elif line.startswith(">"):
            p = text_frame.add_paragraph()
            p.text = line[1:].strip()
            p.level = 0
            p._element.get_or_add_pPr().remove_all("buChar")  # remove bullet if any
            
        else:
            # fallback: treat as paragraph
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.level = 0
            
def presentation2pptx(presentation: SlidesgenPresentation, template_file: str, output_file: str):
    template = Presentation(
        pptx = template_file
    )
        
    template_slides = template.slides
    
    for slide_index, slide in enumerate(presentation.slides):
        try:
            if isinstance(slide, TitleSlide):
                template_layout = template_slides[0].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                set_formatted_content(template.slides[-1].shapes[0], slide.master_title)
                set_formatted_content(template.slides[-1].shapes[1], slide.subtitle)
                      
            elif isinstance(slide, AgendaSlide):
                template_layout = template_slides[1].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                for section_index, section_header in enumerate(slide.sections_header):
                    set_formatted_content(template.slides[-1].shapes[section_index * 2], section_header.id)
                    set_formatted_content(template.slides[-1].shapes[section_index * 2 + 1], section_header.title)
            
            elif isinstance(slide, SectionTransitionSlide):
                template_layout = template_slides[2].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                
                set_formatted_content(template.slides[-1].shapes[0], slide.section_header.id)
                set_formatted_content(template.slides[-1].shapes[1], slide.section_header.title)
                set_formatted_content(template.slides[-1].shapes[2], slide.section_header.subtitle)
                set_formatted_content(template.slides[-1].shapes[3], slide.image)

            elif isinstance(slide, ThankYouSlide):
                template_layout = template_slides[3].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                
                set_formatted_content(template.slides[-1].shapes[0], slide.thank_you_text)
                set_formatted_content(template.slides[-1].shapes[1], slide.additional_info)
                set_formatted_content(template.slides[-1].shapes[2], slide.contact_information)
                
            elif isinstance(slide, SimpleContentSlide):
                template_layout = template_slides[4].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                
                set_formatted_content(template.slides[-1].shapes[0], slide.title)
                set_formatted_content(template.slides[-1].shapes[1], slide.content)
                
            elif isinstance(slide, TwoColumnsSlide):
                template_layout = template_slides[5].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                set_formatted_content(template.slides[-1].shapes[0], slide.title)
                set_formatted_content(template.slides[-1].shapes[1], slide.column1_title)
                set_formatted_content(template.slides[-1].shapes[2], slide.column1_content)
                set_formatted_content(template.slides[-1].shapes[3], slide.column2_title)
                set_formatted_content(template.slides[-1].shapes[4], slide.column2_content)

            elif isinstance(slide, QuoteSlide):
                template_layout = template_slides[6].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                
                set_formatted_content(template.slides[-1].shapes[0], slide.quote)
                set_formatted_content(template.slides[-1].shapes[1], slide.author)
                
            elif isinstance(slide, ImpressionSlide):
                template_layout = template_slides[7].slide_layout
                template.slides.add_slide(template_layout)
                
                template.slides[-1].name = f"{slide_index}"
                
                set_formatted_content(template.slides[-1].shapes[0], slide.impression_text)
                set_formatted_content(template.slides[-1].shapes[1], slide.description)
                
        except Exception as e:
            print(e)
            
    template.save(
        file = output_file
    )
